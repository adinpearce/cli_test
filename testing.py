#!/usr/bin/env python
# -*- coding: UTF-8 -*-
from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches 


#開啟新的簡報物件
prs = Presentation()
#建立簡報檔第一張頁面物件
title_slide_layout = prs.slide_layouts[0] 
#增加一張簡報
slide = prs.slides.add_slide(title_slide_layout)
#設定第一張簡報的標題 
title = slide.shapes.title
title.text = "Hello Python PPT"
#將簡報物件存檔



def add_slide(prs, layout,img):
	slide = prs.slides.add_slide(layout)
	shapes = slide.shapes
	
    #投影片標題
	title_shape = shapes.title
	title_shape.text = '長方圖示範'

	#投影片內第一大點
	body_shape = shapes.placeholders[1]
	tf = body_shape.text_frame
	tf.text = '這是一個一週銷售狀況的示範'

    #投影片內第二階層小點
	p = tf.add_paragraph()
	p.text = '分別有 A、B、C 三種品牌的銷售量'
	p.level = 1


	# show the figure
	left = Inches(3)
	height = Inches(4.5)
	left = top = Inches(3)
	pic = slide.shapes.add_picture(img, left, top, height=height)

	
	return slide


x_labels = ['A','B','C']
sales_num = [120,90,60]
plt.bar(x_labels,sales_num)
plt.savefig('graph.jpg')
img = 'graph.jpg'


title_slide_layout = prs.slide_layouts[1]
slide = add_slide(prs, title_slide_layout ,img)

prs.save("python_ppt_v1.pptx")
